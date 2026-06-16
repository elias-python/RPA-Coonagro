"""
Gera a apresentação PowerPoint do projeto RPA Coonagro.
Executar: python gerar_apresentacao.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paleta de cores ──────────────────────────────────────────────────────────
VERDE_MOSAIC = RGBColor(0x00, 0x6B, 0x3C)  # verde corporativo Mosaic
VERDE_CLARO = RGBColor(0x43, 0xA0, 0x47)
CINZA_ESCURO = RGBColor(0x21, 0x21, 0x21)
CINZA_MEDIO = RGBColor(0x61, 0x61, 0x61)
CINZA_CLARO = RGBColor(0xF5, 0xF5, 0xF5)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
AZUL_SAP = RGBColor(0x00, 0x6E, 0xAF)
LARANJA = RGBColor(0xFF, 0x98, 0x00)
VERMELHO = RGBColor(0xEF, 0x53, 0x50)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # layout em branco


# ── Helpers ──────────────────────────────────────────────────────────────────


def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    l,
    t,
    w,
    h,
    font_size=18,
    bold=False,
    color=CINZA_ESCURO,
    align=PP_ALIGN.LEFT,
    wrap=True,
    italic=False,
):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_box_with_text(
    slide,
    l,
    t,
    w,
    h,
    fill,
    title,
    body,
    title_size=13,
    body_size=11,
    title_color=BRANCO,
    body_color=BRANCO,
    border=None,
):
    shape = add_rect(slide, l, t, w, h, fill, border, 1.5 if border else 0)
    txBox = slide.shapes.add_textbox(
        Inches(l + 0.12), Inches(t + 0.08), Inches(w - 0.24), Inches(h - 0.16)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    # título
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = title_color
    # corpo
    if body:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = body_color
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=CINZA_MEDIO, width_pt=1.5):
    """Seta horizontal simples usando conector."""
    from pptx.util import Emu

    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)  # MSO_CONNECTOR_TYPE.STRAIGHT
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Fundo verde
add_rect(slide, 0, 0, 13.33, 7.5, VERDE_MOSAIC)

# Faixa inferior cinza
add_rect(slide, 0, 6.2, 13.33, 1.3, CINZA_ESCURO)

# Título
add_text(
    slide,
    "RPA — Automação de Lançamento Coonagro",
    0.6,
    1.8,
    12,
    1.2,
    font_size=36,
    bold=True,
    color=BRANCO,
    align=PP_ALIGN.CENTER,
)

# Subtítulo
add_text(
    slide,
    "Sistema automatizado de recebimento de XML fiscal,\nprocessamento SAP e arquivamento de Notas Fiscais de Remessa",
    0.6,
    3.1,
    12,
    1.2,
    font_size=20,
    color=RGBColor(0xC8, 0xE6, 0xC9),
    align=PP_ALIGN.CENTER,
)

# Linha separadora branca
add_rect(slide, 3.5, 4.55, 6.33, 0.04, BRANCO)

# Info rodapé
add_text(
    slide,
    "The Mosaic Company  |  Controladoria PGA1  |  2026",
    0.6,
    6.4,
    12,
    0.6,
    font_size=13,
    color=RGBColor(0xBD, 0xBD, 0xBD),
    align=PP_ALIGN.CENTER,
)

# Ícone decorativo (círculo verde claro)
add_rect(slide, 5.9, 4.8, 1.53, 1.0, VERDE_CLARO)
add_text(
    slide,
    "⚙ RPA",
    5.95,
    4.88,
    1.4,
    0.8,
    font_size=22,
    bold=True,
    color=BRANCO,
    align=PP_ALIGN.CENTER,
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEMA / CONTEXTO
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, VERDE_MOSAIC)
add_text(
    slide,
    "Contexto e Problema",
    0.4,
    0.12,
    12,
    0.7,
    font_size=26,
    bold=True,
    color=BRANCO,
)

# Coluna esquerda — situação anterior
add_rect(slide, 0.4, 1.1, 5.8, 5.9, RGBColor(0xFF, 0xEB, 0xEB))
add_text(
    slide,
    "⚠  Processo Manual (antes)",
    0.55,
    1.2,
    5.5,
    0.5,
    font_size=16,
    bold=True,
    color=VERMELHO,
)

dores = [
    "• Recebimento de XML por e-mail → download manual",
    "• Conferência individual de cada nota fiscal",
    "• Lançamento item a item no SAP (AUTCARR + NFR)",
    "• Impressão manual de Nota de Remessa (J1BNFE)",
    "• Sem controle de quais NFs já foram lançadas",
    "• Alto risco de erro humano e retrabalho",
    "• Processo lento: horas por lote de NFs",
]
y = 1.75
for d in dores:
    add_text(slide, d, 0.6, y, 5.5, 0.42, font_size=12, color=CINZA_ESCURO)
    y += 0.42

# Coluna direita — solução
add_rect(slide, 6.8, 1.1, 6.1, 5.9, RGBColor(0xE8, 0xF5, 0xE9))
add_text(
    slide,
    "✔  Processo Automatizado (agora)",
    6.95,
    1.2,
    5.8,
    0.5,
    font_size=16,
    bold=True,
    color=VERDE_CLARO,
)

ganhos = [
    "• Monitoramento contínuo de XMLs no Outlook",
    "• Filtragem automática CFOP 5902 / 5124",
    "• Lançamento SAP automático (AUTCARR → NFR)",
    "• Impressão automática via J1BNFE quando necessário",
    "• Controle de NFs lançadas (nfs_lancadas.txt)",
    "• Exportação XLSX + PDF para arquivamento",
    "• Dashboard de status em tempo real no Excel",
]
y = 1.75
for g in ganhos:
    add_text(slide, g, 6.95, y, 5.8, 0.42, font_size=12, color=CINZA_ESCURO)
    y += 0.42

# Seta central
add_text(
    slide,
    "→",
    6.1,
    3.7,
    0.7,
    0.6,
    font_size=36,
    bold=True,
    color=VERDE_MOSAIC,
    align=PP_ALIGN.CENTER,
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARQUITETURA DO SISTEMA
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, VERDE_MOSAIC)
add_text(
    slide,
    "Arquitetura do Sistema",
    0.4,
    0.12,
    12,
    0.7,
    font_size=26,
    bold=True,
    color=BRANCO,
)

# Camadas
layers = [
    (
        0.3,
        1.1,
        2.6,
        1.2,
        AZUL_SAP,
        "📧  Outlook",
        "XMLs fiscais recebidos\nda Coonagro",
    ),
    (
        3.3,
        1.1,
        2.6,
        1.2,
        AZUL_SAP,
        "🐍  XML Monitoring.py",
        "Monitor contínuo 24h\nDetecta novos XMLs",
    ),
    (
        6.3,
        1.1,
        2.6,
        1.2,
        AZUL_SAP,
        "🗃  SQLite (.db)",
        "Armazena itens\nfiltrados por CFOP",
    ),
    (
        9.3,
        1.1,
        3.5,
        1.2,
        AZUL_SAP,
        "📊  Excel (.xlsm)",
        "Painel de controle\ne status visual",
    ),
]
for l, t, w, h, fill, title, body in layers:
    add_box_with_text(slide, l, t, w, h, fill, title, body, title_size=13, body_size=10)

# Setas camada 1
for x in [2.9, 5.9, 8.9]:
    add_text(
        slide,
        "→",
        x,
        1.45,
        0.5,
        0.4,
        font_size=22,
        bold=True,
        color=AZUL_SAP,
        align=PP_ALIGN.CENTER,
    )

# Camada SAP
add_box_with_text(
    slide,
    0.3,
    3.0,
    5.6,
    1.3,
    VERDE_MOSAIC,
    "🖥  SAP PA1 — ZT_MM_94N",
    "Passagem 1: AUTCARR  →  Passagem 2: NFR\nImpressão automática J1BNFE quando necessário",
    title_size=14,
    body_size=11,
)

add_box_with_text(
    slide,
    6.3,
    3.0,
    6.5,
    1.3,
    VERDE_CLARO,
    "📁  Arquivamento automático",
    "Exportação XLSX + PDF  |  Pasta Exportacoes\nnfs_lancadas.txt  →  controle de reprocessamento",
    title_size=14,
    body_size=11,
)

# Seta Excel → SAP
add_text(
    slide,
    "▼",
    10.05,
    2.35,
    0.5,
    0.5,
    font_size=18,
    bold=True,
    color=VERDE_MOSAIC,
    align=PP_ALIGN.CENTER,
)
add_text(
    slide,
    "Macro VBA\nIniciarComControle",
    9.5,
    2.4,
    1.7,
    0.55,
    font_size=9,
    color=CINZA_MEDIO,
    align=PP_ALIGN.CENTER,
)

# Seta Monitoring → SAP
add_text(
    slide,
    "▼",
    1.55,
    2.35,
    0.5,
    0.5,
    font_size=18,
    bold=True,
    color=AZUL_SAP,
    align=PP_ALIGN.CENTER,
)

# Seta SAP → Arquivamento
add_text(
    slide,
    "→",
    5.85,
    3.5,
    0.6,
    0.4,
    font_size=22,
    bold=True,
    color=VERDE_MOSAIC,
    align=PP_ALIGN.CENTER,
)

# Legenda de status
add_text(
    slide,
    "Status visual na planilha:",
    0.3,
    4.6,
    4,
    0.4,
    font_size=12,
    bold=True,
    color=CINZA_ESCURO,
)
status_items = [
    (RGBColor(0x42, 0xA5, 0xF5), "AZUL — AUTCARR_OK (aguardando NFR)"),
    (RGBColor(0x66, 0xBB, 0x6A), "VERDE — Concluído"),
    (RGBColor(0xEF, 0x53, 0x50), "VERMELHO — Erro"),
    (RGBColor(0xFF, 0x98, 0x00), "LARANJA — Atenção / Reprocessar"),
]
x_off = 0.3
for cor, label in status_items:
    add_rect(slide, x_off, 5.15, 0.28, 0.28, cor)
    add_text(
        slide, label, x_off + 0.35, 5.12, 2.9, 0.35, font_size=10, color=CINZA_ESCURO
    )
    x_off += 3.22


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FLUXO DE PROCESSAMENTO
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, VERDE_MOSAIC)
add_text(
    slide,
    "Fluxo de Processamento",
    0.4,
    0.12,
    12,
    0.7,
    font_size=26,
    bold=True,
    color=BRANCO,
)

steps = [
    (
        AZUL_SAP,
        "1",
        "Detecção XML",
        "Monitor lê Outlook\nFiltro CFOP 5902/5124\nGrava no SQLite",
    ),
    (
        AZUL_SAP,
        "2",
        "Atualização Excel",
        "Python → Excel via COM\nChave de Acesso em texto\nDashboard atualizado",
    ),
    (
        VERDE_MOSAIC,
        "3",
        "Cockpit SAP",
        "Gera nfs_para_processar.txt\nPython preenche pedidos\nno banco de dados",
    ),
    (
        VERDE_MOSAIC,
        "4",
        "AUTCARR (Pass 1)",
        "SAP ZT_MM_94N\nGrava Pedido + Data\nPeso BAL = MENGE",
    ),
    (
        VERDE_MOSAIC,
        "5",
        "NFR (Pass 2)",
        "SAP ZT_MM_94N\nNF Origem + Série + Qtd\nConfirma e salva",
    ),
    (
        LARANJA,
        "6",
        "J1BNFE (se necessário)",
        "NF não impressa detectada\nExtração automática NF/Série\nImprime e retenta (RETRY)",
    ),
    (
        VERDE_CLARO,
        "7",
        "Exportação",
        "XLSX + PDF em Exportacoes/\nGrava nfs_lancadas.txt\nRemove linhas VERDE",
    ),
]

box_w = 1.6
box_h = 2.2
gap = 0.18
x_start = 0.25

for i, (fill, num, title, body) in enumerate(steps):
    x = x_start + i * (box_w + gap)
    # Número do passo
    add_rect(slide, x, 1.05, box_w, 0.45, fill)
    add_text(
        slide,
        f"Passo {num}",
        x,
        1.08,
        box_w,
        0.4,
        font_size=13,
        bold=True,
        color=BRANCO,
        align=PP_ALIGN.CENTER,
    )
    # Caixa
    add_rect(slide, x, 1.5, box_w, box_h, CINZA_CLARO, border_color=fill, border_pt=2)
    add_text(
        slide,
        title,
        x + 0.05,
        1.55,
        box_w - 0.1,
        0.45,
        font_size=12,
        bold=True,
        color=fill,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        body,
        x + 0.08,
        2.05,
        box_w - 0.16,
        1.6,
        font_size=10,
        color=CINZA_ESCURO,
        align=PP_ALIGN.CENTER,
    )
    # Seta entre passos
    if i < len(steps) - 1:
        ax = x + box_w + 0.01
        add_text(
            slide,
            "›",
            ax,
            2.2,
            gap + 0.14,
            0.5,
            font_size=20,
            bold=True,
            color=CINZA_MEDIO,
            align=PP_ALIGN.CENTER,
        )

# Nota de retry
add_rect(slide, 0.25, 3.85, 12.83, 0.55, RGBColor(0xFF, 0xF9, 0xC4))
add_text(
    slide,
    "♻  Reprocessamento automático: grupos LARANJA (atenção) e AZUL (AUTCARR_OK pendente) são retomados no próximo ciclo de Iniciar.",
    0.4,
    3.9,
    12.5,
    0.45,
    font_size=11,
    color=RGBColor(0x79, 0x55, 0x00),
)

# Nota Recarregar
add_rect(slide, 0.25, 4.55, 12.83, 0.55, RGBColor(0xE8, 0xEA, 0xF6))
add_text(
    slide,
    "🔄  Botão Recarregar: busca todas as NFs do banco (sem filtro de data) excluindo as já lançadas — ideal para importar XMLs históricos.",
    0.4,
    4.6,
    12.5,
    0.45,
    font_size=11,
    color=RGBColor(0x1A, 0x23, 0x7E),
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — COMPONENTES TÉCNICOS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, VERDE_MOSAIC)
add_text(
    slide,
    "Componentes Técnicos",
    0.4,
    0.12,
    12,
    0.7,
    font_size=26,
    bold=True,
    color=BRANCO,
)

componentes = [
    (
        AZUL_SAP,
        "🐍  Python 3.14",
        [
            "XML Monitoring.py — monitor principal 24h",
            "forcar_atualizacao.py — Recarregar manual",
            "force_cockpit.py — prepara lote SAP",
            "openpyxl + pandas + pywin32 + sqlite3",
        ],
    ),
    (
        VERDE_MOSAIC,
        "📊  VBA (.xlsm)",
        [
            "modDashboard.bas — painel de controle",
            "modSAP_Processamento.bas — automação SAP",
            "modPainel.bas — utilitários e exportação",
            "Auto_Open, shapes xlFreeFloating",
        ],
    ),
    (
        RGBColor(0x5C, 0x6B, 0xC0),
        "🗃  Dados",
        [
            "SQLite: dados_rpa_coonagro.db",
            "nfs_lancadas.txt — controle de lançados",
            "Exportacoes/ — XLSX + PDF arquivados",
            "SharePoint (OneDrive) como repositório",
        ],
    ),
    (
        RGBColor(0x00, 0x83, 0x8A),
        "🖥  SAP GUI Scripting",
        [
            "ZT_MM_94N — lançamento AUTCARR + NFR",
            "J1BNFE — impressão Nota de Remessa",
            "Variante NMELO4 pré-configurada",
            "3 níveis de validação de conexão",
        ],
    ),
]

col_w = 2.9
for i, (fill, title, items) in enumerate(componentes):
    col = i % 2
    row = i // 2
    x = 0.35 + col * (col_w + 3.65)
    y = 1.1 + row * 2.8

    add_rect(slide, x, y, col_w + 3.3, 0.5, fill)
    add_text(
        slide,
        title,
        x + 0.1,
        y + 0.08,
        col_w + 3.1,
        0.38,
        font_size=14,
        bold=True,
        color=BRANCO,
    )
    add_rect(
        slide,
        x,
        y + 0.5,
        col_w + 3.3,
        2.1,
        CINZA_CLARO,
        border_color=fill,
        border_pt=1.5,
    )
    yi = y + 0.65
    for item in items:
        add_text(
            slide,
            f"• {item}",
            x + 0.15,
            yi,
            col_w + 3.0,
            0.42,
            font_size=11,
            color=CINZA_ESCURO,
        )
        yi += 0.44


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PRÓXIMOS PASSOS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 0.9, VERDE_MOSAIC)
add_text(
    slide, "Próximos Passos", 0.4, 0.12, 12, 0.7, font_size=26, bold=True, color=BRANCO
)

proximos = [
    (
        AZUL_SAP,
        "🚀  Auto-disparo",
        "Após detecção de novos XMLs, o sistema acionará automaticamente o Cockpit e o Iniciar — sem nenhuma intervenção manual.",
    ),
    (
        LARANJA,
        "📧  E-mail de Falha",
        "Em caso de erro SAP ou NF problemática, o sistema enviará e-mail de alerta automático com detalhes do grupo falhado.",
    ),
    (
        VERDE_MOSAIC,
        "🖥  VM Dedicada (Azure)",
        "Servidor Windows na nuvem Azure com SAP PA1 logado permanentemente, XML Monitoring rodando 24h sem depender da máquina do usuário.",
    ),
    (
        RGBColor(0x5C, 0x6B, 0xC0),
        "📦  GitHub (código)",
        "Repositório github.com/elias-python/RPA-Coonagro já ativo. Na VM: git clone para deploy imediato do código mais recente.",
    ),
]

y = 1.15
for fill, title, desc in proximos:
    add_rect(slide, 0.4, y, 0.08, 1.1, fill)
    add_rect(slide, 0.55, y, 12.38, 0.45, fill)
    add_text(
        slide, title, 0.65, y + 0.06, 11.8, 0.38, font_size=14, bold=True, color=BRANCO
    )
    add_rect(
        slide, 0.55, y + 0.45, 12.38, 0.65, CINZA_CLARO, border_color=fill, border_pt=1
    )
    add_text(slide, desc, 0.7, y + 0.5, 12.1, 0.58, font_size=12, color=CINZA_ESCURO)
    y += 1.3


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ENCERRAMENTO
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, VERDE_MOSAIC)
add_rect(slide, 0, 5.8, 13.33, 1.7, CINZA_ESCURO)

add_text(
    slide,
    "Obrigado!",
    0.6,
    1.8,
    12,
    1.0,
    font_size=48,
    bold=True,
    color=BRANCO,
    align=PP_ALIGN.CENTER,
)

add_text(
    slide,
    "RPA Coonagro  —  The Mosaic Company",
    0.6,
    3.0,
    12,
    0.6,
    font_size=20,
    color=RGBColor(0xC8, 0xE6, 0xC9),
    align=PP_ALIGN.CENTER,
)

add_rect(slide, 4.5, 3.85, 4.33, 0.04, BRANCO)

add_text(
    slide,
    "Controladoria PGA1  |  2026",
    0.6,
    4.1,
    12,
    0.5,
    font_size=14,
    color=RGBColor(0xA5, 0xD6, 0xA7),
    align=PP_ALIGN.CENTER,
)

add_text(
    slide,
    "github.com/elias-python/RPA-Coonagro",
    0.6,
    6.05,
    12,
    0.5,
    font_size=13,
    color=RGBColor(0xBD, 0xBD, 0xBD),
    align=PP_ALIGN.CENTER,
)


# ── Salva ────────────────────────────────────────────────────────────────────
saida = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "Apresentacao_RPA_Coonagro.pptx"
)
prs.save(saida)
print(f"[OK] Apresentação salva em:\n     {saida}")
